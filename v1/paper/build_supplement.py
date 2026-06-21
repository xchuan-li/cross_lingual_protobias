#!/usr/bin/env python3
"""Build the supplement slides that fill gaps in ProgressReport_1.pptx.

Two outputs:
  (standalone)  python code/experiments/build_supplement.py
                -> SUPPLEMENT_additions.pptx  (cover + 5 review slides)
  (merge)       python code/experiments/build_supplement.py --merge
                -> ProgressReport_1_merged.pptx
                   = the full talk with the SC/PA illustration inserted after the
                     Method slide, and the 4 backup slides appended to the appendix.

Content slides (data pulled from real repo files: summary.csv,
exp1_attribute_breakdown.csv, translations.json):
  - Real SC vs PA example images (Animals + Objects only; demography cropped out)
  - Translation sanity check (en/zh/ar/hi)
  - Full socio_attr x language table
  - object x Hindi = .553 (the one low-resource signal) + caveat
  - Confounds (adversarial knob) + ethical handling of demography
"""
import os
import sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
ASSET = os.path.join(ROOT, "code/experiments/assets/protobias_examples_animals_objects.png")
ASSET_RATIO = 2.062  # w/h of the cropped Animals+Objects figure
MAIN = os.path.join(ROOT, "ProgressReport_1.pptx")
OUT_STANDALONE = os.path.join(ROOT, "SUPPLEMENT_additions.pptx")
OUT_MERGED = os.path.join(ROOT, "ProgressReport_1_merged.pptx")

# palette / fonts lifted from ProgressReport_1.pptx
INK   = RGBColor(0x1F, 0x29, 0x33)
GRAY  = RGBColor(0x7B, 0x84, 0x94)
RED   = RGBColor(0xC4, 0x4E, 0x52)
BLUE  = RGBColor(0x4C, 0x72, 0xB0)
LINE  = RGBColor(0xD9, 0xDC, 0xE0)
PANEL = RGBColor(0xF5, 0xF3, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TITLEF, BODYF = "Georgia", "Calibri"

EMU_IN = 914400
SW, SH = 12192000, 6858000
MARGIN = int(0.7 * EMU_IN)


def _set(run, size, color, font=BODYF, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic


def box(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, first=False):
    return tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()


def rect(slide, x, y, w, h, fill, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w or 0.75)
    return sp


def header(slide, kicker, title):
    tf = box(slide, SW - int(6.5 * EMU_IN), int(0.32 * EMU_IN), int(6.0 * EMU_IN), int(0.3 * EMU_IN))
    p = para(tf, True); p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = "Multilingual Prototypicality Bias  ·  Progress Report"
    _set(r, 9, GRAY)
    rect(slide, MARGIN, int(0.62 * EMU_IN), int(0.16 * EMU_IN), int(0.16 * EMU_IN), RED)
    tf3 = box(slide, MARGIN + int(0.26 * EMU_IN), int(0.60 * EMU_IN), int(9 * EMU_IN), int(0.3 * EMU_IN))
    r = para(tf3, True).add_run(); r.text = kicker; _set(r, 11, RED, bold=True)
    tf4 = box(slide, MARGIN, int(0.92 * EMU_IN), int(12 * EMU_IN), int(0.7 * EMU_IN))
    r = para(tf4, True).add_run(); r.text = title; _set(r, 30, INK, font=TITLEF, bold=True)


def bullet(tf, text, first=False, color=INK, size=15, bold=False, gap=6):
    p = para(tf, first); p.space_after = Pt(gap)
    r = p.add_run(); r.text = ("•  " + text) if not bold else text
    _set(r, size, color, bold=bold)
    return p


# ---------------------------------------------------------------- slide builders
def s_cover(slide):
    rect(slide, 0, 0, SW, SH, WHITE)
    rect(slide, MARGIN, int(2.0 * EMU_IN), int(0.16 * EMU_IN), int(2.6 * EMU_IN), RED)
    tf = box(slide, MARGIN + int(0.45 * EMU_IN), int(1.9 * EMU_IN), int(11 * EMU_IN), int(1.4 * EMU_IN))
    r = para(tf, True).add_run(); r.text = "Supplement — additions to the deck"
    _set(r, 40, INK, font=TITLEF, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(8)
    r = p.add_run(); r.text = "Five slides that fill the gaps in ProgressReport_1.pptx."
    _set(r, 16, GRAY)
    tf = box(slide, MARGIN + int(0.45 * EMU_IN), int(3.7 * EMU_IN), int(11 * EMU_IN), int(3 * EMU_IN))
    items = [
        ("1 · Real SC vs PA examples", "the missing illustration — a vision project with no images"),
        ("2 · Translation sanity check", "one item in en / zh / ar / hi — answers the inevitable QA"),
        ("3 · Full socio_attr × language table", "attribute-specific & language-consistent in one view"),
        ("4 · The one low-resource signal", "object × Hindi = .553 — motivates the Exp2 low-resource test"),
        ("5 · Confounds & ethics", "the adversarial-knob confound + handling of demography images"),
    ]
    first = True
    for h, sub in items:
        p = para(tf, first); first = False; p.space_after = Pt(10)
        r = p.add_run(); r.text = h + "   "; _set(r, 16, INK, bold=True)
        r = p.add_run(); r.text = sub; _set(r, 14, GRAY, italic=True)


def s_illustration(slide):
    rect(slide, 0, 0, SW, SH, WHITE)
    header(slide, "ILLUSTRATION  ·  what the judge sees", "SC vs PA: the same trial, two images")
    iw = int(7.4 * EMU_IN); ih = int(iw / ASSET_RATIO)
    ix = MARGIN; iy = int(2.4 * EMU_IN)
    rect(slide, ix - 9000, iy - 9000, iw + 18000, ih + 18000, None, line=LINE, line_w=1)
    slide.shapes.add_picture(ASSET, Emu(ix), Emu(iy), Emu(iw))
    nx = ix + iw + int(0.5 * EMU_IN)
    tf = box(slide, nx, iy, SW - nx - MARGIN, ih, MSO_ANCHOR.TOP)
    bullet(tf, "Each row is one trial: a Semantically-Correct image (SC) vs a Prototypical-but-wrong Adversarial image (PA).", True, size=15)
    bullet(tf, "The prompt stays neutral — it never names the object, only the scene.", size=15)
    bullet(tf, "Check marks = which automatic metrics still prefer the WRONG-but-typical PA image.", size=15)
    bullet(tf, "Demography pairs exist too, but stay as labels / numbers — never shown.", size=14, color=GRAY)
    p = tf.add_paragraph(); p.space_before = Pt(10)
    r = p.add_run(); r.text = "Source: ProtoBias (Roy, Bhatia & Eger, 2026), Fig. 1 — Animals & Objects rows."
    _set(r, 11, GRAY, italic=True)


def s_translation(slide):
    rect(slide, 0, 0, SW, SH, WHITE)
    header(slide, "BACKUP  ·  QA: “how good are the translations?”", "Translation sanity check — one item, four languages")
    rows = [
        ("EN  (source)", "A bird sits on the ground with a small stone in the background near a large log.", False),
        ("ZH  中文", "一只鸟坐在地上，背景中靠近一根大圆木的地方有一块小石头。", False),
        ("AR  العربية", "طائر يجلس على الأرض مع حجر صغير في الخلفية بالقرب من جذع شجرة كبير.", True),
        ("HI  हिन्दी", "एक पक्षी एक बड़े लट्ठे के पास पृष्ठभूमि में एक छोटे पत्थर के साथ जमीन पर बैठा है।", True),
    ]
    ty = int(2.05 * EMU_IN); rh = int(0.92 * EMU_IN); lw = int(1.7 * EMU_IN)
    for i, (lang, txt, rtl) in enumerate(rows):
        yy = ty + i * (rh + int(0.12 * EMU_IN))
        rect(slide, MARGIN, yy, SW - 2 * MARGIN, rh, PANEL if i % 2 == 0 else WHITE, line=LINE, line_w=0.75)
        tf = box(slide, MARGIN + int(0.2 * EMU_IN), yy, lw, rh, MSO_ANCHOR.MIDDLE)
        r = para(tf, True).add_run(); r.text = lang; _set(r, 14, RED if i else BLUE, bold=True)
        tf = box(slide, MARGIN + lw + int(0.2 * EMU_IN), yy, SW - 2 * MARGIN - lw - int(0.5 * EMU_IN), rh, MSO_ANCHOR.MIDDLE)
        p = para(tf, True); p.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
        r = p.add_run(); r.text = txt; _set(r, 15, INK)
    tf = box(slide, MARGIN, ty + 4 * (rh + int(0.12 * EMU_IN)) + int(0.05 * EMU_IN), SW - 2 * MARGIN, int(0.7 * EMU_IN))
    r = para(tf, True).add_run()
    r.text = "Backend: Google (deep-translator). ZH human-verified (native). AR/HI back-translation check queued for Exp2 (see Limitations)."
    _set(r, 13, GRAY, italic=True)


def s_attr_table(slide):
    rect(slide, 0, 0, SW, SH, WHITE)
    header(slide, "BACKUP  ·  full breakdown", "Full socio_attr × language error rates")
    data = [
        ("Wealth",    ".774", ".774", ".839", ".710", "31", True),
        ("Power",     ".703", ".672", ".594", ".719", "64", True),
        ("Civility",  ".547", ".484", ".547", ".578", "64", False),
        ("Morality",  ".473", ".559", ".505", ".409", "93", False),
        ("Intellect", ".417", ".479", ".542", ".479", "48", False),
    ]
    cols = ["socio_attr", "en", "zh", "ar", "hi", "n / lang"]
    tx = MARGIN; ty = int(2.0 * EMU_IN)
    cw = [int(2.6 * EMU_IN)] + [int(1.7 * EMU_IN)] * 4 + [int(1.6 * EMU_IN)]
    rh = int(0.62 * EMU_IN)
    cx = tx
    for j, c in enumerate(cols):
        rect(slide, cx, ty, cw[j], rh, INK)
        tf = box(slide, cx, ty, cw[j], rh, MSO_ANCHOR.MIDDLE)
        p = para(tf, True); p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = c; _set(r, 14, WHITE, bold=True)
        tf.margin_left = Emu(int(0.12 * EMU_IN)); cx += cw[j]
    for i, row in enumerate(data):
        yy = ty + (i + 1) * rh
        name, en, zh, ar, hi, n, strong = row
        cx = tx
        for j, v in enumerate([name, en, zh, ar, hi, n]):
            rect(slide, cx, yy, cw[j], rh, PANEL if strong else WHITE, line=LINE, line_w=0.75)
            tf = box(slide, cx, yy, cw[j], rh, MSO_ANCHOR.MIDDLE)
            p = para(tf, True); p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = v
            col = RED if (j in (1, 2, 3, 4) and strong) else INK
            _set(r, 14, col, bold=(j == 0) or strong)
            tf.margin_left = Emu(int(0.12 * EMU_IN)); cx += cw[j]
    tf = box(slide, tx, ty + 6 * rh + int(0.15 * EMU_IN), SW - 2 * MARGIN, int(1.2 * EMU_IN))
    bullet(tf, "Wealth & Power (shaded) sit well above chance in every language; Morality & Intellect hover at .50 everywhere.", True, size=15)
    bullet(tf, "The attribute ordering is preserved across en/zh/ar/hi — the effect is the attribute, not the language.", size=15, color=BLUE)
    tf2 = box(slide, tx, ty + 6 * rh + int(0.15 * EMU_IN) + int(0.95 * EMU_IN), SW - 2 * MARGIN, int(0.4 * EMU_IN))
    r = para(tf2, True).add_run(); r.text = "Source: exp2_attribute_bias/exp1_attribute_breakdown.csv  ·  chance = 0.50"
    _set(r, 11, GRAY, italic=True)


def s_object_hindi(slide):
    rect(slide, 0, 0, SW, SH, WHITE)
    header(slide, "BACKUP  ·  motivates Outlook step 3", "The one low-resource signal: object × Hindi")
    data = [
        ("animal",     ".523", ".575", ".570", ".573", None),
        ("demography", ".560", ".577", ".573", ".553", None),
        ("object",     ".467", ".470", ".463", ".553", "hi"),
    ]
    cols = ["domain", "en", "zh", "ar", "hi"]
    tx = MARGIN; ty = int(2.0 * EMU_IN)
    cw = [int(2.6 * EMU_IN)] + [int(1.5 * EMU_IN)] * 4
    rh = int(0.6 * EMU_IN)
    cx = tx
    for j, c in enumerate(cols):
        rect(slide, cx, ty, cw[j], rh, INK)
        tf = box(slide, cx, ty, cw[j], rh, MSO_ANCHOR.MIDDLE)
        p = para(tf, True); p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = c; _set(r, 14, WHITE, bold=True)
        tf.margin_left = Emu(int(0.12 * EMU_IN)); cx += cw[j]
    for i, row in enumerate(data):
        yy = ty + (i + 1) * rh
        vals = [row[0]] + list(row[1:5]); flag = row[5]
        cx = tx
        for j, v in enumerate(vals):
            hot = (flag == "hi" and j == 4)
            rect(slide, cx, yy, cw[j], rh, RED if hot else WHITE, line=LINE, line_w=0.75)
            tf = box(slide, cx, yy, cw[j], rh, MSO_ANCHOR.MIDDLE)
            p = para(tf, True); p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = v
            _set(r, 14, WHITE if hot else INK, bold=(j == 0) or hot)
            tf.margin_left = Emu(int(0.12 * EMU_IN)); cx += cw[j]
    nx = tx + cw[0] + 4 * cw[1] + int(0.6 * EMU_IN)
    tf = box(slide, nx, ty, SW - nx - MARGIN, int(2.6 * EMU_IN))
    bullet(tf, "On Objects, en/zh/ar all read .46–.47 (no bias) — but Hindi alone jumps to .553.", True, size=15)
    bullet(tf, "It is the only cell where the low-resource language behaves differently.", size=15)
    bullet(tf, "Locally consistent with “low-resource = more biased” — the hypothesis Exp2 will test directly.", size=15, color=BLUE)
    cy = ty + 4 * rh + int(0.55 * EMU_IN)
    rect(slide, tx, cy, SW - 2 * MARGIN, int(1.0 * EMU_IN), PANEL, line=RED, line_w=1)
    tf = box(slide, tx + int(0.25 * EMU_IN), cy, SW - 2 * MARGIN - int(0.5 * EMU_IN), int(1.0 * EMU_IN), MSO_ANCHOR.MIDDLE)
    r = para(tf, True).add_run()
    r.text = "Do NOT over-claim: n = 300, 95% CI [.497, .609] still touches 0.50. One cell, one model, one seed — a lead to confirm, not a result."
    _set(r, 14, INK, bold=True)


def s_confound_ethics(slide):
    rect(slide, 0, 0, SW, SH, WHITE)
    header(slide, "HONESTY  ·  add to Limitations", "Confounds we own & how we handle the data")
    colw = int(5.5 * EMU_IN)
    lx = MARGIN; ly = int(2.05 * EMU_IN); ch = int(3.6 * EMU_IN)
    rect(slide, lx, ly, colw, ch, WHITE, line=LINE, line_w=1)
    rect(slide, lx, ly, colw, int(0.55 * EMU_IN), INK)
    tf = box(slide, lx + int(0.25 * EMU_IN), ly, colw, int(0.55 * EMU_IN), MSO_ANCHOR.MIDDLE)
    r = para(tf, True).add_run(); r.text = "The adversarial-knob confound"; _set(r, 16, WHITE, font=TITLEF, bold=True)
    tf = box(slide, lx + int(0.25 * EMU_IN), ly + int(0.75 * EMU_IN), colw - int(0.5 * EMU_IN), ch - int(0.9 * EMU_IN))
    bullet(tf, "The PA image is built by turning ONE visual knob: count, colour/tone, layout, spatial, or scale.", True, size=14)
    bullet(tf, "So a wrong pick may track that low-level change, not 'prototypicality' per se.", size=14)
    bullet(tf, "We flag it openly; Exp2's mixed-effects model adds knob as a control term.", size=14, color=BLUE)
    rx = lx + colw + int(0.5 * EMU_IN)
    rect(slide, rx, ly, colw, ch, WHITE, line=LINE, line_w=1)
    rect(slide, rx, ly, colw, int(0.55 * EMU_IN), RED)
    tf = box(slide, rx + int(0.25 * EMU_IN), ly, colw, int(0.55 * EMU_IN), MSO_ANCHOR.MIDDLE)
    r = para(tf, True).add_run(); r.text = "Handling the demography data"; _set(r, 16, WHITE, font=TITLEF, bold=True)
    tf = box(slide, rx + int(0.25 * EMU_IN), ly + int(0.75 * EMU_IN), colw - int(0.5 * EMU_IN), ch - int(0.9 * EMU_IN))
    bullet(tf, "500 demography pairs encode race / religion / orientation stereotypes — sensitive material.", True, size=14)
    bullet(tf, "Public slides use only Animal / Object images; demography stays as labels & numbers.", size=14)
    bullet(tf, "We study the stereotype to measure it, never to endorse it — stated up front.", size=14, color=RED)
    tf = box(slide, MARGIN, ly + ch + int(0.3 * EMU_IN), SW - 2 * MARGIN, int(0.6 * EMU_IN))
    r = para(tf, True).add_run(); r.text = "Naming these first turns the obvious QA attacks into points in our favour."
    _set(r, 14, GRAY, italic=True)


def s_subcategory(slide):
    rect(slide, 0, 0, SW, SH, WHITE)
    header(slide, "BACKUP  ·  drilling into demography",
           "Sub-population: the judge defaults to the unmarked image")
    # left table: subcategory rate / n / CI
    tx = MARGIN; ty = int(2.0 * EMU_IN)
    cols = ["sub-population", "error", "n", "95% CI"]
    cw = [int(2.5 * EMU_IN), int(1.0 * EMU_IN), int(0.8 * EMU_IN), int(1.6 * EMU_IN)]
    rh = int(0.6 * EMU_IN)
    data = [
        ("Sexual orientation", ".748", "416", "[.70, .79]", True),
        ("Nationality",        ".474", "392", "[.43, .52]", False),
        ("Religion",           ".464", "392", "[.42, .51]", False),
    ]
    cx = tx
    for j, c in enumerate(cols):
        rect(slide, cx, ty, cw[j], rh, INK)
        tf = box(slide, cx, ty, cw[j], rh, MSO_ANCHOR.MIDDLE)
        p = para(tf, True); p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = c; _set(r, 13, WHITE, bold=True)
        tf.margin_left = Emu(int(0.12 * EMU_IN)); cx += cw[j]
    for i, row in enumerate(data):
        yy = ty + (i + 1) * rh
        name, er, n, ci, strong = row
        cx = tx
        for j, v in enumerate([name, er, n, ci]):
            rect(slide, cx, yy, cw[j], rh, PANEL if strong else WHITE, line=LINE, line_w=0.75)
            tf = box(slide, cx, yy, cw[j], rh, MSO_ANCHOR.MIDDLE)
            p = para(tf, True); p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = v
            _set(r, 13, RED if (strong and j in (1, 3)) else INK, bold=(j == 0) or strong)
            tf.margin_left = Emu(int(0.12 * EMU_IN)); cx += cw[j]
    # caption under table: what "error" means here
    tw = sum(cw)
    cap = box(slide, tx, ty + 4 * rh + int(0.12 * EMU_IN), tw, int(0.7 * EMU_IN))
    p = para(cap, True)
    r = p.add_run(); r.text = "error = how often the judge picks the unmarked “default” image — which is the wrong one. "
    _set(r, 12, GRAY, italic=True)
    p2 = cap.add_paragraph()
    r = p2.add_run(); r.text = "chance = .50"; _set(r, 12, GRAY, italic=True)
    # right column: the mechanism (precise, non-overclaiming)
    nx = tx + tw + int(0.55 * EMU_IN)
    tf = box(slide, nx, ty, SW - nx - MARGIN, int(3.3 * EMU_IN))
    bullet(tf, "Mechanism: the CORRECT image is the one showing a marked identity — a pride flag (in 141/169 items), a kippah; the trap is the plain, unmarked “default” person.", True, size=14)
    bullet(tf, "On sexual-orientation items the judge picks that unmarked default ~75% of the time — i.e. it rejects the correct marked image.", size=14, color=RED)
    bullet(tf, "Religion & nationality ≈ chance. The same effect inflates the Wealth headline (Wealth × sexual-orientation = .97).", size=14)
    # bottom band: salience confound + Exp2 + ethics
    by = int(5.55 * EMU_IN)
    rect(slide, MARGIN, by, SW - 2 * MARGIN, int(1.25 * EMU_IN), PANEL, line=RED, line_w=1)
    tf = box(slide, MARGIN + int(0.25 * EMU_IN), by, SW - 2 * MARGIN - int(0.5 * EMU_IN), int(1.25 * EMU_IN), MSO_ANCHOR.MIDDLE)
    p = para(tf, True)
    r = p.add_run(); r.text = "Confound — salience vs identity: "; _set(r, 13, RED, bold=True)
    r = p.add_run(); r.text = "the pride markers are visually loud, so we can’t yet tell “avoids the LGBTQ-coded image” from “avoids the busier image.”"; _set(r, 13, INK)
    p2 = tf.add_paragraph(); p2.space_before = Pt(4)
    r = p2.add_run(); r.text = "Exp2 fix: "; _set(r, 13, RED, bold=True)
    r = p2.add_run(); r.text = "picked_adv ~ socio_attr + subcategory + knob + lang + gender + (1|item) to disentangle.   "; _set(r, 13, INK)
    r = p2.add_run(); r.text = "We measure the model’s bias, never endorse it — numbers only."; _set(r, 12, GRAY, italic=True)


BACKUPS = [s_translation, s_attr_table, s_object_hindi, s_confound_ethics]

# Speaker note for the subcategory slide (EN / 中文 format, matches SPEAKER_NOTES.md)
SUBCAT_NOTE = (
    "EN\n"
    "Let me be precise here, because this cut is easy to overclaim. In these demography "
    "items the CORRECT image is the one showing a marked identity — a pride flag, a kippah "
    "— and the trap image is the plain, unmarked 'default' person. On the sexual-orientation "
    "items the model picks that unmarked default about 75% of the time, even though it's the "
    "wrong answer; religion and nationality are at chance. And this same effect is what "
    "inflates the wealth number — wealth-by-sexual-orientation is .97. The honest caveat: "
    "those pride markers are visually loud, so we can't yet separate 'the model avoids "
    "LGBTQ-coded images' from 'the model avoids the busier image.' Exp2 adds sub-population "
    "and the visual knob to the model to pull those apart. And to be clear — we're measuring "
    "the model's bias, not endorsing any stereotype, and we only show numbers.\n\n"
    "中文\n"
    "这里我要讲准一点，因为这个切片很容易被过度解读。在这些 demography 样本里，“语义正确”的"
    "那张图恰好是带可见少数群体标记的（彩虹旗、犹太基帕），而“典型陷阱”图是那个没有任何标记的"
    "“默认人”。在性取向这一类上，模型大约 75% 的时候选了那个无标记默认人，哪怕它是错的答案；"
    "宗教和国籍则在随机线上。而且正是这个效应把财富那个数顶高了——财富 × 性取向那格是 .97。"
    "要诚实说的 caveat：那些彩虹标记在视觉上很显眼，所以我们目前分不开“模型在回避 LGBTQ 相关"
    "的图”和“模型只是在回避更花哨的图”。Exp2 会把人群类型和视觉旋钮一起放进模型来拆开这两者。"
    "还有，我们是在测量模型的偏见，不是认同任何刻板印象，而且只展示数字。"
)


def add_slide(prs, layout, builder):
    slide = prs.slides.add_slide(layout)
    builder(slide)
    return slide


def move_slide(prs, from_idx, to_idx):
    lst = prs.slides._sldIdLst
    els = list(lst)
    el = els[from_idx]
    lst.remove(el)
    lst.insert(to_idx, el)


def build_standalone():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    blank = prs.slide_layouts[6]
    for b in [s_cover, s_illustration] + BACKUPS:
        add_slide(prs, blank, b)
    prs.save(OUT_STANDALONE)
    print("wrote", OUT_STANDALONE, "—", len(list(prs.slides._sldIdLst)), "slides")


def build_merged():
    prs = Presentation(MAIN)
    layout = prs.slide_layouts[0]  # DEFAULT (blank, 0 placeholders)
    n0 = len(list(prs.slides._sldIdLst))  # 18
    # 1) illustration -> insert right after Method (slide 6, idx 5) => becomes slide 7
    add_slide(prs, layout, s_illustration)         # lands at end (idx n0)
    move_slide(prs, n0, 6)                          # move to index 6
    # 2) the 4 backups -> append to the very end (after existing appendix)
    for b in BACKUPS:
        add_slide(prs, layout, b)
    prs.save(OUT_MERGED)
    print("wrote", OUT_MERGED, "—", len(list(prs.slides._sldIdLst)), "slides")


def _find_slide_index(prs, needle):
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and needle in sh.text_frame.text:
                return i
    return None


def _notes_template(prs):
    for s in prs.slides:
        for ph in s.notes_slide.placeholders:
            if ph.placeholder_format.idx == 1:
                return ph._element
    return None


def patch_subcat(path, with_note):
    from copy import deepcopy
    prs = Presentation(path)
    layout = prs.slide_layouts[0]
    template_sp = _notes_template(prs) if with_note else None
    add_slide(prs, layout, s_subcategory)                      # appended at end
    new_idx = len(list(prs.slides._sldIdLst)) - 1
    anchor = _find_slide_index(prs, "Full socio_attr")         # the socio_attr table slide
    pos = (anchor + 1) if anchor is not None else new_idx
    move_slide(prs, new_idx, pos)
    if with_note:
        target = list(prs.slides)[pos]
        ns = target.notes_slide
        if ns.notes_text_frame is None and template_sp is not None:
            ns.shapes._spTree.append(deepcopy(template_sp))
        ns.notes_text_frame.text = SUBCAT_NOTE
    prs.save(path)
    n = len(list(prs.slides._sldIdLst))
    print(f"patched {path} -> {n} slides (subcat at slide {pos + 1}{', with note' if with_note else ''})")


def patch_both():
    patch_subcat(OUT_MERGED, with_note=False)
    patch_subcat(os.path.join(ROOT, "ProgressReport_1_merged (Note).pptx"), with_note=True)


if __name__ == "__main__":
    if "--merge" in sys.argv:
        build_merged()
    elif "--add-subcat" in sys.argv:
        patch_both()
    else:
        build_standalone()
